# -*- coding: utf-8 -*-
"""
轻言OPS PPT 生成脚本
风格：深色科技风 · AI 产品发布会视觉
主色：#050B18 / #0A1F44 / #00C8FF / #7B61FF
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os
import sys

# ============================================================
# 配色
# ============================================================
BG_DEEP = (0x05, 0x0B, 0x18)       # 黑曜石黑 背景
BG_CARD = (0x0A, 0x1F, 0x44)       # 深空蓝 卡片
ACCENT_BLUE = (0x00, 0xC8, 0xFF)   # 科技蓝
ACCENT_PURPLE = (0x7B, 0x61, 0xFF) # AI 紫
WHITE = (0xFF, 0xFF, 0xFF)
GRAY = (0x94, 0xA3, 0xB8)
SILVER = (0xCB, 0xD5, 0xE1)
ORANGE = (0xF9, 0x73, 0x16)
GREEN = (0x10, 0xB9, 0x81)
RED = (0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUTPUT = os.path.join(os.path.dirname(os.path.dirname(__file__)), '轻言OPS_演示PPT.pptx')


def rgb(t):
    from pptx.dml.color import RGBColor  # 延迟导入确保环境干净
    return RGBColor(*t)

def R(shp):
    return shp.left + shp.width

def B(shp):
    return shp.top + shp.height


def set_bg(slide, color_tuple=BG_DEEP):
    """设置幻灯片纯色背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color_tuple)
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_rect(slide, x, y, w, h, fill=BG_CARD, line=None, line_w=0.75,
             radius=False, alpha=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    if radius:
        shp.adjustments[0] = 0.08
    return shp


def add_line(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width=1.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 直线
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(width)
    return ln


def add_text(slide, x, y, w, h, text, size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name='微软雅黑'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = font_name
    # 中文字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': font_name})
    rPr.append(ea)
    return tb


def add_bullet(slide, x, y, w, h, lines, size=13, color=SILVER, line_sp=1.3,
               font_name='微软雅黑'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    first = True
    for item in lines:
        if isinstance(item, tuple):
            t, c, b = item
        else:
            t, c, b = item, color, False
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.line_spacing = Pt(int(size * line_sp))
        run = p.add_run()
        run.text = '· ' + t
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(c)
        run.font.bold = b
        run.font.name = font_name
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {'typeface': font_name})
        rPr.append(ea)
    return tb


def add_header(slide, title, subtitle=None, section_tag=None):
    """每页顶部：章节标签 + 大标题 + 副标题"""
    # 顶部科技蓝线
    add_rect(slide, 0, 0, SLIDE_W, Emu(38100), fill=ACCENT_BLUE)
    # 章节小标签
    if section_tag:
        tag = add_rect(slide, Inches(0.6), Inches(0.55), Inches(2.2), Inches(0.36),
                       fill=BG_CARD, line=ACCENT_BLUE, line_w=0.75, radius=True)
        add_text(slide, tag.left, tag.top + Emu(30000), tag.width, Emu(250000),
                 '  ' + section_tag, size=11, color=ACCENT_BLUE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)
    # 主标题
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.75),
             title, size=28, bold=True, color=WHITE)
    # 辅助装饰线
    add_line(slide, Inches(0.6), Inches(1.75), Inches(2.8), Inches(1.75),
             color=ACCENT_PURPLE, width=2)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GRAY)


def slide_number(slide, idx):
    add_text(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
             f'{idx:02d}/12', size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ============================================================
# 各页内容
# ============================================================

def slide_01_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
    set_bg(s)
    # 背景装饰：深空蓝斜条
    deco = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, SLIDE_W, SLIDE_H)
    deco.fill.solid()
    deco.fill.fore_color.rgb = rgb((0x08, 0x14, 0x2E))
    deco.rotation = 180
    deco.line.fill.background()
    s.shapes._spTree.remove(deco._element)
    s.shapes._spTree.insert(3, deco._element)

    # 中心发光圆（数字人核心示意）
    core_x, core_y = Inches(9.2), Inches(3.2)
    for r, col in [(Inches(2.4), (0x0A, 0x1F, 0x44)),
                   (Inches(1.6), (0x0F, 0x2B, 0x5E)),
                   (Inches(0.9), (0x1A, 0x3A, 0x7A))]:
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, core_x - r, core_y - r, r * 2, r * 2)
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(col)
        shp.line.color.rgb = rgb(ACCENT_BLUE)
        shp.line.width = Pt(1.0)
    center = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                core_x - Inches(0.5), core_y - Inches(0.5),
                                Inches(1), Inches(1))
    center.fill.solid()
    center.fill.fore_color.rgb = rgb(ACCENT_BLUE)
    center.line.fill.background()
    # 中心文字
    add_text(s, core_x - Inches(0.5), core_y - Inches(0.18), Inches(1), Inches(0.4),
             'AI', size=22, bold=True, color=BG_DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 连接线
    for ang in (0, 60, 120, 180, 240, 300):
        import math
        rad = math.radians(ang)
        x1 = core_x + Inches(1.0) * math.cos(rad)
        y1 = core_y + Inches(1.0) * math.sin(rad)
        x2 = core_x + Inches(2.2) * math.cos(rad)
        y2 = core_y + Inches(2.2) * math.sin(rad)
        add_line(s, x1, y1, x2, y2, color=ACCENT_PURPLE, width=1.2)
        # 节点
        nd = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                x2 - Inches(0.08), y2 - Inches(0.08),
                                Inches(0.16), Inches(0.16))
        nd.fill.solid()
        nd.fill.fore_color.rgb = rgb(ACCENT_PURPLE)
        nd.line.fill.background()

    # 左侧标题
    add_text(s, Inches(0.8), Inches(2.0), Inches(8), Inches(1.6),
             '轻言OPS', size=72, bold=True, color=WHITE)
    # 标题下划线
    add_line(s, Inches(0.8), Inches(3.55), Inches(4.5), Inches(3.55),
             color=ACCENT_BLUE, width=3)
    add_text(s, Inches(0.8), Inches(3.7), Inches(8), Inches(0.6),
             '大模型驱动的数字人智能运维助手', size=24, color=SILVER)
    add_text(s, Inches(0.8), Inches(4.3), Inches(8), Inches(0.5),
             '面向 IT 运维场景 · 全离线运行 · 一键部署', size=16, color=GRAY)

    # 三大价值色块
    vals = [
        ('智能', '大模型 · 语音零门槛', ACCENT_BLUE),
        ('安全', '全离线 · 二次确认', GREEN),
        ('便捷', '一键启动 · 一站式', ORANGE),
    ]
    bx = Inches(0.8)
    by = Inches(5.4)
    bw = Inches(2.3)
    bh = Inches(1.2)
    for i, (t, d, c) in enumerate(vals):
        card = add_rect(s, bx + i * (bw + Inches(0.15)), by, bw, bh,
                        fill=(0x08, 0x14, 0x2E), line=c, line_w=1, radius=True)
        add_text(s, card.left + Inches(0.2), card.top + Inches(0.1),
                 Inches(2), Inches(0.5), t, size=18, bold=True, color=c)
        add_text(s, card.left + Inches(0.2), card.top + Inches(0.62),
                 Inches(2), Inches(0.5), d, size=11, color=GRAY)

    # 底部信息
    add_text(s, Inches(0.8), Inches(7.0), Inches(7), Inches(0.3),
             '闫吉乐  ·  2026 年 8 月', size=11, color=GRAY)
    add_text(s, Inches(8.5), Inches(7.0), Inches(4.3), Inches(0.3),
             'github.com/YANJI-AFK/qingyan-OPS', size=11, color=ACCENT_BLUE,
             align=PP_ALIGN.RIGHT)


def slide_02_contents(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, 'C O N T E N T S', '12 章节 · 从痛点到落地的完整叙事', '00 · 目录')
    slide_number(s, 2)

    # 12 章节分两列排布（神经网络路径风格卡片）
    left_items = [
        ('01', '行业痛点', '运维效率的"最后一公里"'),
        ('02', '产品定位', '三合一语音运维平台'),
        ('03', '全离线 AI 技术栈', '4 大独立 AI 引擎'),
        ('04', '意图解析 + FSM', '双重策略 · 5 状态护栏'),
        ('05', '语音交互全链路', '三引擎 TTS · 双层 VAD'),
        ('06', '技术架构总览', '前后端分离 4 层架构'),
    ]
    right_items = [
        ('07', '功能演示 ①', '工单管理 + 统计看板'),
        ('08', '功能演示 ②', '3D 监控大屏 + 实时 KPI'),
        ('09', '功能演示 ③', '排班管理 + 语音排班'),
        ('10', '一键部署', '13 步全自动启动'),
        ('11', '未来路线图', 'V1.0 → V3.0 三阶段'),
        ('12', '总结与 Q&A', '张口就来 · AI 运维'),
    ]

    def draw_col(items, start_x, col_color):
        y = Inches(2.5)
        for i, (num, title, desc) in enumerate(items):
            card = add_rect(s, start_x, y, Inches(5.8), Inches(0.72),
                            fill=BG_CARD, line=col_color, line_w=0.75, radius=True)
            # 编号圆
            circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                      card.left + Inches(0.2), card.top + Inches(0.13),
                                      Inches(0.46), Inches(0.46))
            circ.fill.solid()
            circ.fill.fore_color.rgb = rgb(col_color)
            circ.line.fill.background()
            add_text(s, circ.left, circ.top, circ.width, circ.height,
                     num, size=13, bold=True, color=BG_DEEP,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 标题
            add_text(s, card.left + Inches(0.82), card.top + Inches(0.1),
                     Inches(2.6), Inches(0.3), title,
                     size=15, bold=True, color=WHITE)
            add_text(s, card.left + Inches(0.82), card.top + Inches(0.42),
                     Inches(4.8), Inches(0.3), desc,
                     size=10.5, color=GRAY)
            # 连接线
            if i < len(items) - 1:
                add_line(s, start_x + Inches(0.43), y + Inches(0.78),
                         start_x + Inches(0.43), y + Inches(0.78) + Emu(210000),
                         color=col_color, width=1.2)
            y += Inches(0.78 + 0.06)

    draw_col(left_items, Inches(0.6), ACCENT_BLUE)
    draw_col(right_items, Inches(6.95), ACCENT_PURPLE)


def slide_03_pain(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '运维人员 70% 时间花在"找信息"上',
               '传统运维入口分散 · 操作繁琐 · 安全风险高', '01 · 行业痛点')
    slide_number(s, 3)

    # 左侧示意：多屏幕堆叠
    sc = add_rect(s, Inches(0.6), Inches(2.5), Inches(5.4), Inches(4.4),
                  fill=(0x08, 0x14, 0x2E), radius=True, line=ACCENT_BLUE, line_w=0.75)
    # 屏幕标题栏
    bar = add_rect(s, sc.left + Inches(0.05), sc.top + Inches(0.05),
                   sc.width - Inches(0.1), Inches(0.4), fill=BG_CARD, radius=True)
    # 三盏灯
    for j, c in enumerate([RED, ORANGE, GREEN]):
        lx = bar.left + Inches(0.15) + j * Inches(0.3)
        lamp = s.shapes.add_shape(MSO_SHAPE.OVAL, lx, bar.top + Inches(0.12),
                                  Inches(0.16), Inches(0.16))
        lamp.fill.solid()
        lamp.fill.fore_color.rgb = rgb(c)
        lamp.line.fill.background()
    add_text(s, bar.left + Inches(1.2), bar.top + Inches(0.08),
             Inches(4), Inches(0.3), '多系统切换 · 运维工程师桌面',
             size=11, color=SILVER)

    # 三个堆叠小窗口
    wx = sc.left + Inches(0.3)
    wy = sc.top + Inches(0.7)
    for k, (wt, col) in enumerate([('工单系统', ACCENT_BLUE),
                                    ('监控平台', ACCENT_PURPLE),
                                    ('排班系统', ORANGE)]):
        win = add_rect(s, wx + k * Inches(0.5), wy + k * Inches(0.5),
                       Inches(3.8), Inches(1.0),
                       fill=BG_CARD, line=col, line_w=1, radius=True)
        add_text(s, win.left + Inches(0.18), win.top + Inches(0.2),
                 Inches(3), Inches(0.3), wt, size=13, bold=True, color=col)
        # 乱码色块表示内容
        for kk in range(3):
            add_rect(s, win.left + Inches(0.18) + (kk % 2) * Inches(1.6),
                     win.top + Inches(0.55) + (kk // 2) * Inches(0.22),
                     Inches(1.4), Inches(0.13), fill=(0x11, 0x25, 0x50))
    # 红色警告线
    add_line(s, sc.left + Inches(0.3), sc.top + Inches(3.8),
             sc.left + Inches(5.0), sc.top + Inches(3.8), color=RED, width=2)
    add_text(s, sc.left + Inches(0.3), sc.top + Inches(3.9),
             Inches(5), Inches(0.3),
             '⚠  一次故障排查需切换 4 次页面 · 人均耗时 15 分钟',
             size=11, bold=True, color=RED)

    # 右侧 4 个痛点卡片
    pains = [
        ('01', '多系统切换低效', '工单/监控/排班分散在 3~5 个系统', '夜班不便打字，仅支持键鼠操作', ACCENT_BLUE),
        ('02', '语音交互缺失', '82% 夜间运维希望语音操作', '传统系统无语音输入/输出能力', ACCENT_PURPLE),
        ('03', '操作缺乏护栏', '0.3% 失误带来 30% 额外工作量', '删除/排班/指派无二次确认', ORANGE),
        ('04', '数据安全风险', '云端 AI 数据出境合规风险', '断网即不可用 · 成本高昂', RED),
    ]
    cx = Inches(6.4)
    cy = Inches(2.5)
    cw = Inches(6.4)
    ch = Inches(1.02)
    for i, (n, t, d1, d2, c) in enumerate(pains):
        card = add_rect(s, cx, cy + i * (ch + Inches(0.05)),
                        cw, ch, fill=BG_CARD, line=c, line_w=0.75, radius=True)
        # 编号
        add_text(s, card.left + Inches(0.25), card.top + Inches(0.12),
                 Inches(0.7), Inches(0.4), n, size=22, bold=True, color=c)
        add_line(s, card.left + Inches(1.0), card.top + Inches(0.18),
                 card.left + Inches(1.0), card.top + card.height - Inches(0.2),
                 color=c, width=1.5)
        add_text(s, card.left + Inches(1.18), card.top + Inches(0.1),
                 Inches(5), Inches(0.3), t, size=14, bold=True, color=WHITE)
        add_text(s, card.left + Inches(1.18), card.top + Inches(0.42),
                 Inches(5), Inches(0.24), '- ' + d1, size=10, color=SILVER)
        add_text(s, card.left + Inches(1.18), card.top + Inches(0.67),
                 Inches(5), Inches(0.24), '- ' + d2, size=10, color=GRAY)


def slide_04_position(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '轻言OPS = 语音对话 + 大模型 + 运维平台',
               '三合一 · 让运维人员"说一句话"就完成查询与操作', '02 · 产品定位')
    slide_number(s, 4)

    # 中心三环（太阳系结构）
    cx, cy = Inches(4.6), Inches(4.8)
    rings = [
        (Inches(2.7), ACCENT_BLUE,  '内层'),
        (Inches(1.9), ACCENT_PURPLE,'中层'),
        (Inches(1.1), ORANGE,       '核心'),
    ]
    # 最外层环
    out = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             cx - rings[0][0], cy - rings[0][0],
                             rings[0][0] * 2, rings[0][0] * 2)
    out.fill.background()
    out.line.color.rgb = rgb(rings[0][1])
    out.line.width = Pt(1.5)
    # 中间
    mid = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             cx - rings[1][0], cy - rings[1][0],
                             rings[1][0] * 2, rings[1][0] * 2)
    mid.fill.background()
    mid.line.color.rgb = rgb(rings[1][1])
    mid.line.width = Pt(1.5)
    # 核心
    core = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              cx - rings[2][0], cy - rings[2][0],
                              rings[2][0] * 2, rings[2][0] * 2)
    core.fill.solid()
    core.fill.fore_color.rgb = rgb((0x1A, 0x3A, 0x7A))
    core.line.color.rgb = rgb(rings[2][1])
    core.line.width = Pt(2)
    add_text(s, core.left, core.top + Inches(0.3),
             core.width, Inches(0.5), '数字人', size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, core.left, core.top + Inches(0.85),
             core.width, Inches(0.35), 'AI 核心', size=11, color=ORANGE,
             align=PP_ALIGN.CENTER, bold=True)

    # 内层标签：语音交互
    tags_inner = [
        ('ASR 识别', -80, ACCENT_BLUE), ('TTS 合成', -30, ACCENT_BLUE),
        ('双层 VAD', 45, ACCENT_BLUE),
    ]
    import math
    for txt, ang_deg, c in tags_inner:
        rad = math.radians(ang_deg)
        px = cx + rings[1][0] * 1.0 * math.cos(rad)
        py = cy + rings[1][0] * 1.0 * math.sin(rad)
        tg = add_rect(s, px - Inches(0.9), py - Inches(0.2),
                      Inches(1.8), Inches(0.4), fill=BG_CARD,
                      line=c, line_w=0.75, radius=True)
        add_text(s, tg.left, tg.top + Inches(0.07), tg.width, Inches(0.28),
                 txt, size=10.5, bold=True, color=c, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 中层标签：大模型决策
    tags_mid = [
        ('Ollama + qwen3:8b', -110, ACCENT_PURPLE),
        ('11 类意图解析', 100, ACCENT_PURPLE),
        ('6 层 JSON 容错', 160, ACCENT_PURPLE),
        ('二次确认护栏', 0, ACCENT_PURPLE),
    ]
    for txt, ang_deg, c in tags_mid:
        rad = math.radians(ang_deg)
        px = cx + rings[0][0] * 1.02 * math.cos(rad)
        py = cy + rings[0][0] * 1.02 * math.sin(rad)
        tg = add_rect(s, px - Inches(1.15), py - Inches(0.22),
                      Inches(2.3), Inches(0.44), fill=(0x08, 0x14, 0x2E),
                      line=c, line_w=0.75, radius=True)
        add_text(s, tg.left, tg.top + Inches(0.08), tg.width, Inches(0.28),
                 txt, size=11, bold=True, color=c, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 外层标签：业务执行（右列表形式）
    add_text(s, Inches(8.6), Inches(2.5), Inches(4.4), Inches(0.4),
             '业务执行层', size=15, bold=True, color=WHITE)
    add_line(s, Inches(8.6), Inches(2.95), Inches(12.8), Inches(2.95),
             color=GREEN, width=2)
    biz = [
        ('工单管理', 'CRUD · SLA · 流转 · 智能指派', GREEN),
        ('监控大屏', '3D 地球 · 实时 KPI · 拓扑', GREEN),
        ('排班管理', '日历视图 · 语音排班 · 批量修改', GREEN),
        ('人员管理', '档案 · 岗位 · 标签 · 搜索', GREEN),
        ('统计看板', '趋势 · 分布 · 负荷 · SLA', GREEN),
        ('系统监控', 'psutil 实时 CPU/内存/磁盘', GREEN),
    ]
    ly = Inches(3.15)
    for i, (t, d, c) in enumerate(biz):
        y = ly + i * Inches(0.68)
        # 指示线
        import math as mm
        ang = mm.radians(-20 + i * 4)
        x1 = cx + rings[0][0] * mm.cos(ang)
        y1 = cy + rings[0][0] * mm.sin(ang)
        add_line(s, x1, y1, Inches(8.6), y + Inches(0.16), color=GREEN, width=0.8)
        # 卡片
        card = add_rect(s, Inches(8.6), y, Inches(4.4), Inches(0.58),
                        fill=BG_CARD, line=c, line_w=0.75, radius=True)
        add_text(s, card.left + Inches(0.2), card.top + Inches(0.08),
                 Inches(2), Inches(0.25), t, size=12.5, bold=True, color=c)
        add_text(s, card.left + Inches(0.2), card.top + Inches(0.32),
                 Inches(4), Inches(0.2), d, size=9.5, color=GRAY)

    # 底部金句
    qbox = add_rect(s, Inches(0.6), Inches(6.9), Inches(12.2), Inches(0.45),
                    fill=(0x08, 0x14, 0x2E), line=ACCENT_BLUE, line_w=0.75, radius=True)
    add_text(s, qbox.left + Inches(0.3), qbox.top + Inches(0.1),
             Inches(11.5), Inches(0.3),
             '金句：说一句话，完成一次运维操作。',
             size=13, bold=True, color=ACCENT_BLUE)


def slide_05_ai_stack(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '核心创新 ① · 全离线 AI 技术栈',
               '4 大独立 AI 引擎协同工作 · 全链路本地运行 · 零外部 API 依赖',
               '03 · 全离线 AI')
    slide_number(s, 5)

    cards = [
        ('LLM Brain', 'Ollama + qwen3:8b (Q4 量化)',
         '决策大脑', '5.2 GB', '~1s', ACCENT_PURPLE),
        ('AI Ear', 'FunASR paraformer-zh + VAD',
         '耳朵', '~1 GB', '~200ms', GREEN),
        ('AI Voice', 'sherpa / SAPI / edge-tts',
         '嘴巴', '~345 MB (3 模型)', '~300ms', ACCENT_BLUE),
        ('AI Memory', 'SQL Server + pyodbc (8 连接池)',
         '记忆', '-', '<50ms', ORANGE),
    ]
    # 左上 2 + 右下 2 网格
    positions = [(Inches(0.6), Inches(2.5)),
                 (Inches(6.95), Inches(2.5)),
                 (Inches(0.6), Inches(4.5)),
                 (Inches(6.95), Inches(4.5))]
    for (name, tech, role, size, lat, c), (x, y) in zip(cards, positions):
        card = add_rect(s, x, y, Inches(6.15), Inches(1.85),
                        fill=BG_CARD, line=c, line_w=1, radius=True)
        # 角标角色色标
        add_rect(s, card.left, card.top, Inches(0.12), card.height, fill=c, radius=False)
        add_text(s, card.left + Inches(0.4), card.top + Inches(0.2),
                 Inches(3.5), Inches(0.32), name, size=18, bold=True, color=c)
        add_text(s, card.left + Inches(0.4), card.top + Inches(0.62),
                 Inches(5.6), Inches(0.32), tech, size=13, color=WHITE)
        # 参数
        add_text(s, card.left + Inches(0.4), card.top + Inches(1.1),
                 Inches(1.8), Inches(0.25), f'角色 · {role}',
                 size=10, color=GRAY)
        add_text(s, card.left + Inches(2.4), card.top + Inches(1.1),
                 Inches(1.8), Inches(0.25), f'模型大小 · {size}',
                 size=10, color=GRAY)
        add_text(s, card.left + Inches(4.3), card.top + Inches(1.1),
                 Inches(1.8), Inches(0.25), f'延迟 · {lat}',
                 size=10, color=GRAY)
        add_line(s, card.left + Inches(0.4), card.top + Inches(1.45),
                 R(card) - Inches(0.2), card.top + Inches(1.45),
                 color=c, width=0.75)

    # 中心流程协作图（底部卡片）
    flow = add_rect(s, Inches(0.6), Inches(6.48), Inches(12.5), Inches(0.72),
                    fill=(0x08, 0x14, 0x2E), line=ACCENT_BLUE, line_w=0.75, radius=True)
    steps = ['用户语音', 'ASR 识别', '数字归一化', 'LLM 意图', '业务执行', 'TTS 合成', '语音回复']
    # 横向 7 节点连线
    n = len(steps)
    total_w = flow.width - Inches(0.6)
    gap = total_w / (n - 1)
    xs = [flow.left + Inches(0.3) + gap * i for i in range(n)]
    for i, (st, x) in enumerate(zip(steps, xs)):
        # 节点圆
        nd = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                x - Inches(0.12), flow.top + Inches(0.22),
                                Inches(0.24), Inches(0.24))
        nd.fill.solid()
        nd.fill.fore_color.rgb = rgb(ACCENT_BLUE if i % 2 == 0 else ACCENT_PURPLE)
        nd.line.fill.background()
        add_text(s, x - Inches(1.2), flow.top + Inches(0.02), Inches(2.4), Inches(0.22),
                 st, size=10.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_line(s, x + Inches(0.12), flow.top + Inches(0.34),
                     xs[i + 1] - Inches(0.12), flow.top + Inches(0.34),
                     color=ACCENT_BLUE, width=1.5)


def slide_06_fsm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '核心创新 ② · 双重意图解析 + FSM 安全护栏',
               '11 类意图 · 双重策略 · 5 状态机 = 精准且安全',
               '04 · 意图解析 + FSM')
    slide_number(s, 6)

    # 左侧：双路径意图解析
    add_text(s, Inches(0.6), Inches(2.4), Inches(6), Inches(0.4),
             '双重意图解析策略', size=16, bold=True, color=WHITE)
    add_line(s, Inches(0.6), Inches(2.85), Inches(4), Inches(2.85),
             color=ACCENT_BLUE, width=2)

    # 输入块
    inp = add_rect(s, Inches(0.6), Inches(3.05), Inches(5.9), Inches(0.7),
                   fill=(0x08, 0x14, 0x2E), line=SILVER, line_w=0.75, radius=True)
    add_text(s, inp.left, inp.top + Inches(0.18), inp.width, Inches(0.35),
             '🎐 用户语音输入  →  "让王强排明天早班"',
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 分支线
    branch_y = Inches(3.8)
    add_line(s, Inches(3.55), B(inp), Inches(2.15), branch_y + Inches(0.2),
             color=ACCENT_BLUE, width=1.2)
    add_line(s, Inches(3.55), B(inp), Inches(5.0), branch_y + Inches(0.2),
             color=ACCENT_PURPLE, width=1.2)

    # 路径 1
    p1 = add_rect(s, Inches(0.6), branch_y, Inches(2.9), Inches(1.7),
                  fill=BG_CARD, line=ACCENT_BLUE, line_w=1, radius=True)
    add_text(s, p1.left + Inches(0.2), p1.top + Inches(0.12),
             Inches(2.5), Inches(0.32), '路径 1 · 规则预检',
             size=12, bold=True, color=ACCENT_BLUE)
    add_bullet(s, p1.left + Inches(0.15), p1.top + Inches(0.5),
               Inches(2.7), Inches(1.1),
               ['< 50ms 快速响应', '绕过 LLM · 零 token 消耗',
                ('工单 ID 归一化', ACCENT_BLUE, True)],
               size=10.5, color=SILVER)
    # 闪电图标示意
    add_text(s, p1.left + Inches(2.2), p1.top + Inches(0.12),
             Inches(0.6), Inches(0.4), '⚡', size=20, color=ACCENT_BLUE,
             align=PP_ALIGN.RIGHT)

    # 路径 2
    p2 = add_rect(s, Inches(3.6), branch_y, Inches(2.9), Inches(1.7),
                  fill=BG_CARD, line=ACCENT_PURPLE, line_w=1, radius=True)
    add_text(s, p2.left + Inches(0.2), p2.top + Inches(0.12),
             Inches(2.5), Inches(0.32), '路径 2 · LLM 推理',
             size=12, bold=True, color=ACCENT_PURPLE)
    add_bullet(s, p2.left + Inches(0.15), p2.top + Inches(0.5),
               Inches(2.7), Inches(1.1),
               ['qwen3:8b · temp=0.1', '严格 JSON · max_tokens=500',
                '6 层 JSON 容错 + 降级重试'],
               size=10.5, color=SILVER)
    add_text(s, p2.left + Inches(2.2), p2.top + Inches(0.12),
             Inches(0.6), Inches(0.4), '🧠', size=20, color=ACCENT_PURPLE,
             align=PP_ALIGN.RIGHT)

    # 汇聚
    converge_y = Inches(5.65)
    add_line(s, Inches(2.05), B(p1), Inches(3.55), converge_y,
             color=ACCENT_BLUE, width=1.2)
    add_line(s, Inches(5.05), B(p2), Inches(3.55), converge_y,
             color=ACCENT_PURPLE, width=1.2)

    merge = add_rect(s, Inches(0.6), converge_y, Inches(5.9), Inches(0.7),
                     fill=BG_CARD, line=GREEN, line_w=1, radius=True)
    add_text(s, merge.left, merge.top + Inches(0.18), merge.width, Inches(0.35),
             '✔  归一化意图结果  →  {intent, api, params}',
             size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 右侧：FSM 5 状态
    add_text(s, Inches(7.0), Inches(2.4), Inches(5.9), Inches(0.4),
             'FSM 5 状态安全护栏', size=16, bold=True, color=WHITE)
    add_line(s, Inches(7.0), Inches(2.85), Inches(11.3), Inches(2.85),
             color=ACCENT_PURPLE, width=2)

    states = [
        ('IDLE', '初始 / 操作完成', ACCENT_BLUE, Inches(9.45), Inches(3.05)),
        ('PROCESSING', 'LLM 解析中', GREEN, Inches(7.15), Inches(4.35)),
        ('REPLYING', '生成回复中', ORANGE, Inches(11.75), Inches(4.35)),
        ('AWAITING', '二次确认 · 30s 超时', RED, Inches(7.15), Inches(5.8)),
        ('EXECUTING', 'DB 写操作', ACCENT_PURPLE, Inches(11.75), Inches(5.8)),
    ]
    nodes = {}
    for name, desc, c, x, y in states:
        w, h = Inches(2.2), Inches(0.92)
        node = add_rect(s, x, y, w, h, fill=BG_CARD, line=c, line_w=1, radius=True)
        add_text(s, node.left + Inches(0.12), node.top + Inches(0.1),
                 w - Inches(0.24), Inches(0.35), name,
                 size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, node.left + Inches(0.12), node.top + Inches(0.5),
                 w - Inches(0.24), Inches(0.35), desc,
                 size=9.5, color=GRAY, align=PP_ALIGN.CENTER)
        nodes[name] = (x + w / 2, y + h / 2)

    # 状态转移箭头（FSM）
    transitions = [
        ('IDLE',       'PROCESSING', '接收输入'),
        ('PROCESSING', 'REPLYING',   '查询意图'),
        ('PROCESSING', 'AWAITING',   '写操作意图'),
        ('REPLYING',   'IDLE',       'TTS 完成'),
        ('AWAITING',   'EXECUTING',  '用户确认'),
        ('AWAITING',   'IDLE',       '取消/超时'),
        ('EXECUTING',  'IDLE',       '业务完成'),
    ]
    def arrow_center(sx, sy, ex, ey, c, label='', d=0.1):
        from pptx.enum.shapes import MSO_CONNECTOR
        import math
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)
        conn.line.color.rgb = rgb(c)
        conn.line.width = Pt(1.0)
        # 箭头
        conn.line.end_arrowhead = 1  # 三角
        if label:
            # 垂直偏移中点
            mid_x = (sx + ex) / 2
            mid_y = (sy + ey) / 2 - Inches(0.15)
            add_text(s, mid_x - Inches(1.2), mid_y, Inches(2.4), Inches(0.22),
                     label, size=9, bold=True, color=c, align=PP_ALIGN.CENTER)
    # 简化：直接画主要连线
    connections = [
        # x1,y1, x2,y2, color, label
        (Inches(10.55), Inches(3.45), Inches(8.25), Inches(4.35), ACCENT_BLUE,  '接收输入'),
        (Inches(8.25), Inches(5.27), Inches(8.25), Inches(5.8),  GREEN,         '写操作'),
        (Inches(9.35), Inches(5.27), Inches(11.75), Inches(4.8), ORANGE,        '查询意图'),
        (Inches(12.85), Inches(5.27), Inches(12.85), Inches(3.5), ACCENT_BLUE,  'TTS 完成'),
        (Inches(9.35), Inches(6.26), Inches(11.75), Inches(6.26), ACCENT_PURPLE,'确认 → 执行'),
        (Inches(10.55), Inches(6.72), Inches(10.55), Inches(6.72), RED,         ''),
        (Inches(8.25), Inches(6.72), Inches(10.55), Inches(3.97), RED,          '取消/超时'),
        (Inches(12.85), Inches(6.72), Inches(11.55), Inches(3.97), ACCENT_BLUE, '执行完成'),
    ]
    for x1, y1, x2, y2, c, lab in connections:
        try:
            from pptx.enum.shapes import MSO_CONNECTOR
            cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
            cn.line.color.rgb = rgb(c)
            cn.line.width = Pt(1.0)
            import sys
            if not hasattr(sys, '_skip_arrow'):
                try:
                    cn.line.end_arrowhead_style = 1  # 部分版本兼容
                except Exception:
                    pass
        except Exception:
            add_line(s, x1, y1, x2, y2, c, 1.0)

    # 关键指标（底部）
    kpis = [
        ('3 类操作', '删除/指派/排班 默认二次确认', ORANGE),
        ('30 秒', '超时自动取消 · 防死锁', ACCENT_BLUE),
        ('8 轮', '指代消解 · 连续上下文', ACCENT_PURPLE),
        ('多格式', '工单 ID 归一化容错', GREEN),
    ]
    kx = Inches(0.6)
    ky = Inches(6.5)
    for i, (t, d, c) in enumerate(kpis):
        card = add_rect(s, kx + i * Inches(3.15), ky, Inches(3.0), Inches(0.7),
                        fill=(0x08, 0x14, 0x2E), line=c, line_w=0.75, radius=True)
        add_text(s, card.left + Inches(0.15), card.top + Inches(0.1),
                 Inches(2.7), Inches(0.28), t, size=12, bold=True, color=c)
        add_text(s, card.left + Inches(0.15), card.top + Inches(0.38),
                 Inches(2.7), Inches(0.24), d, size=9.5, color=GRAY)


def slide_07_voice(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '核心创新 ③ · 语音交互全链路 + 三引擎 TTS',
               '从"听到"到"说出" · 端到端语音交互',
               '05 · 语音全链路')
    slide_number(s, 7)

    # 上半：5 节点流水线
    nodes5 = [
        ('🎙 ① 语音采集', 'VAD 自动检测\n手动/连续聆听\n模式自选', ACCENT_BLUE),
        ('🦻 ② ASR 识别', 'FunASR paraformer\nWebM→WAV 16kHz mono\n数字归一化', GREEN),
        ('🧠 ③ LLM 决策', 'Ollama qwen3:8b\n严格 JSON\ntemp=0.1', ACCENT_PURPLE),
        ('🗄 ④ 业务执行', 'Flask + DB\n二次确认护栏\n8 连接池事务', ORANGE),
        ('🔊 ⑤ TTS 输出', 'sherpa→SAPI→edge\n数字按位读\n百分比规范', ACCENT_BLUE),
    ]
    nx = Inches(0.55)
    ny = Inches(2.5)
    nw = Inches(2.42)
    nh = Inches(2.05)
    prev_center = None
    for i, (t, d, c) in enumerate(nodes5):
        x = nx + i * (nw + Inches(0.13))
        card = add_rect(s, x, ny, nw, nh, fill=BG_CARD, line=c, line_w=1, radius=True)
        # 顶栏
        add_rect(s, card.left, card.top, card.width, Inches(0.45), fill=c, radius=True)
        add_text(s, card.left, card.top + Inches(0.1), card.width, Inches(0.3),
                 t, size=12, bold=True, color=BG_DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_bullet(s, card.left + Inches(0.15), card.top + Inches(0.55),
                   nw - Inches(0.3), nh - Inches(0.6),
                   d.split('\n'), size=10.5, color=SILVER)
        # 箭头
        if prev_center:
            add_line(s, prev_center, ny + nh / 2,
                     card.left, ny + nh / 2, color=ACCENT_BLUE, width=2)
        prev_center = R(card)

    # 左下：TTS 三引擎降级
    tts = add_rect(s, Inches(0.6), Inches(4.75), Inches(6.15), Inches(2.5),
                   fill=(0x08, 0x14, 0x2E), line=ACCENT_BLUE, line_w=1, radius=True)
    add_text(s, tts.left + Inches(0.25), tts.top + Inches(0.15),
             Inches(4), Inches(0.35), 'TTS 三引擎降级链',
             size=14, bold=True, color=ACCENT_BLUE)
    chain = [
        ('① sherpa-onnx VITS', '离线 · 神经网络 · 7 音色\n默认：Melo 中英混合女声', GREEN, '推荐'),
        ('② Windows SAPI', '离线兜底 · 自动枚举已安装语音\n晓晓 / 云希 / 晓伊', ORANGE, '降级'),
        ('③ edge-tts', 'TTS_MODE=auto 时启用\n微软在线神经网络', ACCENT_PURPLE, '可选'),
    ]
    yy = tts.top + Inches(0.55)
    for name, d, c, tag in chain:
        box = add_rect(s, tts.left + Inches(0.2), yy,
                       tts.width - Inches(0.4), Inches(0.58),
                       fill=BG_CARD, line=c, line_w=0.75, radius=True)
        add_text(s, box.left + Inches(0.15), box.top + Inches(0.08),
                 Inches(3.5), Inches(0.28), name, size=12, bold=True, color=c)
        add_text(s, R(box) - Inches(1.0), box.top + Inches(0.08),
                 Inches(0.8), Inches(0.28), tag, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT)
        add_text(s, box.left + Inches(0.15), box.top + Inches(0.32),
                 Inches(5.5), Inches(0.22), d, size=9.5, color=GRAY)
        yy += Inches(0.66)
        if tag != '可选':
            # 箭头
            add_line(s, box.left + Inches(2.7), B(box),
                     box.left + Inches(2.7), yy, color=ACCENT_BLUE, width=1.2)

    # 右下：VAD + 数字智能播报（双栏卡）
    # 卡片 1：智能数字播报
    nd = add_rect(s, Inches(6.95), Inches(4.75), Inches(3.0), Inches(2.5),
                  fill=(0x08, 0x14, 0x2E), line=ACCENT_PURPLE, line_w=1, radius=True)
    add_text(s, nd.left + Inches(0.2), nd.top + Inches(0.15),
             Inches(2.6), Inches(0.35), 'TTS 智能数字播报',
             size=13, bold=True, color=ACCENT_PURPLE)
    rules = [
        ('1~3 位', '277 个工单', '二百七十七'),
        ('≥4 位', '2026 年', '二零二六'),
        ('% 号', 'CPU 45.2%', '百分之四十五点二'),
    ]
    ry = nd.top + Inches(0.6)
    for cond, ex, how in rules:
        add_line(s, nd.left + Inches(0.2), ry, R(nd) - Inches(0.2), ry,
                 color=(0x15, 0x2A, 0x5A), width=0.75)
        add_text(s, nd.left + Inches(0.2), ry + Inches(0.05),
                 Inches(0.7), Inches(0.22), cond, size=10, bold=True, color=GREEN)
        add_text(s, nd.left + Inches(1.0), ry + Inches(0.05),
                 Inches(1.8), Inches(0.22), f'示例：{ex}',
                 size=9.5, color=SILVER)
        add_text(s, nd.left + Inches(1.0), ry + Inches(0.3),
                 Inches(1.8), Inches(0.22), f'朗读：{how}',
                 size=9.5, color=ACCENT_BLUE)
        ry += Inches(0.6)

    # 卡片 2：双层 VAD
    vad = add_rect(s, Inches(10.08), Inches(4.75), Inches(3.0), Inches(2.5),
                   fill=(0x08, 0x14, 0x2E), line=ORANGE, line_w=1, radius=True)
    add_text(s, vad.left + Inches(0.2), vad.top + Inches(0.15),
             Inches(2.6), Inches(0.35), '双层 VAD 持续聆听',
             size=13, bold=True, color=ORANGE)
    algos = [
        ('前端 1.8s', 'RMS 50ms 采样\n动态噪声基线\n控制录音停止'),
        ('后端 2400ms', 'FunASR fsmn_vad\n精确句尾判断\n识别性能稳定'),
    ]
    vy = vad.top + Inches(0.6)
    for h, d in algos:
        box = add_rect(s, vad.left + Inches(0.15), vy,
                       vad.width - Inches(0.3), Inches(0.8),
                       fill=BG_CARD, line=ORANGE, line_w=0.75, radius=True)
        add_text(s, box.left + Inches(0.12), box.top + Inches(0.08),
                 Inches(1.8), Inches(0.28), h, size=11, bold=True, color=ORANGE)
        add_text(s, box.left + Inches(0.12), box.top + Inches(0.35),
                 Inches(2.6), Inches(0.5), d, size=9, color=GRAY)
        vy += Inches(0.88)
    add_text(s, vad.left + Inches(0.15), vy, Inches(2.7), Inches(0.25),
             '灵敏度 2.5× · 最短 800ms 过滤噪声',
             size=9.5, color=SILVER)


def slide_08_arch(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '技术架构总览',
               '前后端分离 · 4 层服务栈 · 全离线部署',
               '06 · 架构总览')
    slide_number(s, 8)

    # 4 层垂直架构
    layers = [
        ('L1 · 展示层',
         ['Vue 3 + TS + Vite 8 + Pinia',
          'ECharts 6 · Three.js (3D 地球)',
          'App.vue 浮窗 · 8 页面 + 2 子布局'],
         ACCENT_BLUE, Inches(2.0)),
        ('L2 · API 层',
         ['Axios (baseURL 127.0.0.1:5000)',
          '20+ RESTful 端点',
          'Flask-CORS 跨域'],
         GREEN, Inches(3.05)),
        ('L3 · 服务层',
         ['Flask 3.1 (port 5000)',
          'ASR · LLM · TTS 三服务',
          'pyodbc 8 连接池 · psutil 监控'],
         ACCENT_PURPLE, Inches(4.1)),
        ('L4 · 数据层',
         ['SQL Server 2019+',
          '6 域 11 表 · 工单 / 人员 / 标签',
          '监控域 · 系统域'],
         ORANGE, Inches(5.15)),
    ]
    for name, items, c, y in layers:
        x = Inches(0.6)
        w = Inches(8.2)
        h = Inches(0.95)
        card = add_rect(s, x, y, w, h, fill=BG_CARD, line=c, line_w=1.2, radius=True)
        # 左侧色条
        add_rect(s, card.left, card.top, Inches(0.12), card.height, fill=c)
        add_text(s, card.left + Inches(0.3), card.top + Inches(0.1),
                 Inches(3), Inches(0.35), name,
                 size=14, bold=True, color=c)
        # 项
        tx = card.left + Inches(2.8)
        for k, it in enumerate(items):
            col = tx + k * Inches(1.85)
            add_text(s, col, card.top + Inches(0.3),
                     Inches(1.7), Inches(0.6), it,
                     size=10, color=SILVER, anchor=MSO_ANCHOR.TOP)

    # 右侧：基础设施层 + 关键指标
    infra = add_rect(s, Inches(9.0), Inches(2.0), Inches(4.15), Inches(4.2),
                     fill=(0x08, 0x14, 0x2E), line=SILVER, line_w=0.75, radius=True)
    add_text(s, infra.left + Inches(0.25), infra.top + Inches(0.15),
             Inches(3.8), Inches(0.35), '基础设施层',
             size=14, bold=True, color=SILVER)
    infra_items = [
        ('Python 3.12+', 'venv 虚拟环境', ACCENT_BLUE),
        ('Node.js 22.18+', 'npm 前端构建', GREEN),
        ('ffmpeg', 'C:\\ffmpeg\\bin 音频转换', ACCENT_PURPLE),
        ('Windows 10/11', 'SAPI 5.4 系统语音', ORANGE),
        ('Ollama 本地', 'GPU 推理服务', SILVER),
    ]
    yy = infra.top + Inches(0.6)
    for t, d, c in infra_items:
        row = add_rect(s, infra.left + Inches(0.15), yy,
                       infra.width - Inches(0.3), Inches(0.6),
                       fill=BG_CARD, line=c, line_w=0.75, radius=True)
        add_text(s, row.left + Inches(0.12), row.top + Inches(0.08),
                 Inches(2), Inches(0.25), t, size=11, bold=True, color=c)
        add_text(s, row.left + Inches(0.12), row.top + Inches(0.32),
                 Inches(3.8), Inches(0.2), d, size=9, color=GRAY)
        yy += Inches(0.68)

    # 关键指标（底部横条）
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.35),
             '关键指标', size=13, bold=True, color=WHITE)
    metrics = [
        ('后端代码', '≈ 5,000 行 Python', ACCENT_BLUE),
        ('前端代码', '≈ 8,000 行 TS/Vue', ACCENT_PURPLE),
        ('页面路由', '8 页面 · 2 子布局', GREEN),
        ('API 端点', '20+ RESTful', ORANGE),
        ('测试用例', '60+ 覆盖意图/语音/工单', SILVER),
    ]
    mx = Inches(0.6)
    my = Inches(6.65)
    for t, d, c in metrics:
        card = add_rect(s, mx, my, Inches(2.44), Inches(0.7),
                        fill=(0x08, 0x14, 0x2E), line=c, line_w=0.75, radius=True)
        add_text(s, card.left + Inches(0.15), card.top + Inches(0.1),
                 Inches(2.1), Inches(0.25), t, size=10, bold=True, color=c)
        add_text(s, card.left + Inches(0.15), card.top + Inches(0.35),
                 Inches(2.1), Inches(0.25), d, size=12, bold=True, color=WHITE)
        mx += Inches(2.5)


def slide_09_demo12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '功能演示 ①② · 工单 + 监控',
               '全程语音交互 · 实时数据采集',
               '07 · 工单 & 监控')
    slide_number(s, 9)

    # 左：工单管理 + 统计看板
    left = add_rect(s, Inches(0.6), Inches(2.45), Inches(6.15), Inches(4.8),
                    fill=(0x08, 0x14, 0x2E), line=ACCENT_BLUE, line_w=1, radius=True)
    # 顶部标签
    lb = add_rect(s, left.left + Inches(0.2), left.top + Inches(0.15),
                  Inches(1.6), Inches(0.35), fill=ACCENT_BLUE, radius=True)
    add_text(s, lb.left, lb.top + Inches(0.05), lb.width, Inches(0.25),
             '01 工单管理', size=11, bold=True, color=BG_DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 对话气泡
    dialog_y = left.top + Inches(0.62)
    bubbles = [
        ('最近一周有多少工单已完成？', WHITE, False, Inches(0.25)),
        ('共有 277 个已完成工单。', ACCENT_BLUE, True, Inches(1.45)),
        ('其中多少高优先级的？', WHITE, False, Inches(0.25)),
        ('共有 76 个高优先级已完成工单。', ACCENT_BLUE, True, Inches(1.45)),
    ]
    for txt, c, is_ai, off in bubbles:
        box = add_rect(s, left.left + off, dialog_y,
                       left.width - off - Inches(0.4), Inches(0.52),
                       fill=BG_CARD, line=c, line_w=0.75, radius=True)
        add_text(s, box.left + Inches(0.15), box.top + Inches(0.12),
                 box.width - Inches(0.3), Inches(0.28),
                 ('👤 ' if not is_ai else '🤖 ') + txt,
                 size=11, color=c, bold=is_ai)
        dialog_y += Inches(0.6)
    # 工单特性
    add_line(s, left.left + Inches(0.25), dialog_y,
             R(left) - Inches(0.25), dialog_y,
             color=ACCENT_BLUE, width=0.75)
    add_text(s, left.left + Inches(0.25), dialog_y + Inches(0.08),
             Inches(3), Inches(0.25), '功能特性', size=11, bold=True, color=ACCENT_BLUE)
    feats = [
        'CRUD · 语音新建/修改/查询/指派',
        '7 种状态流转 · SLA 48h 超时预警',
        '智能指派 · 工单 ID 多格式归一化',
    ]
    add_bullet(s, left.left + Inches(0.25), dialog_y + Inches(0.38),
               left.width - Inches(0.5), Inches(1.8),
               feats, size=10.5, color=SILVER)

    # 统计看板小窗
    stats = add_rect(s, left.left + Inches(0.25), dialog_y + Inches(2.05),
                     left.width - Inches(0.5), Inches(1.35),
                     fill=BG_CARD, line=ACCENT_PURPLE, line_w=0.75, radius=True)
    add_text(s, stats.left + Inches(0.15), stats.top + Inches(0.08),
             Inches(3), Inches(0.25), '📊 统计看板',
             size=11, bold=True, color=ACCENT_PURPLE)
    # 4 个数据卡片
    kdata = [('277', '总数', ACCENT_BLUE),
             ('76',  '高优', ACCENT_PURPLE),
             ('25%', '完成率', GREEN),
             ('12',  '本周', ORANGE)]
    for i, (v, l, c) in enumerate(kdata):
        bx = stats.left + Inches(0.15) + i * Inches(1.4)
        by = stats.top + Inches(0.42)
        mini = add_rect(s, bx, by, Inches(1.3), Inches(0.78),
                        fill=(0x06, 0x10, 0x25), line=c, line_w=0.75, radius=True)
        add_text(s, mini.left, mini.top + Inches(0.05), mini.width, Inches(0.4),
                 v, size=18, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, mini.left, mini.top + Inches(0.45), mini.width, Inches(0.22),
                 l, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # 右：3D 监控大屏
    right = add_rect(s, Inches(6.95), Inches(2.45), Inches(6.15), Inches(4.8),
                     fill=(0x08, 0x14, 0x2E), line=ACCENT_PURPLE, line_w=1, radius=True)
    lb2 = add_rect(s, right.left + Inches(0.2), right.top + Inches(0.15),
                   Inches(1.6), Inches(0.35), fill=ACCENT_PURPLE, radius=True)
    add_text(s, lb2.left, lb2.top + Inches(0.05), lb2.width, Inches(0.25),
             '02 3D 监控大屏', size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 地球示意圆
    gx, gy = right.left + Inches(2.0), right.top + Inches(2.0)
    earth = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               gx - Inches(1.15), gy - Inches(1.15),
                               Inches(2.3), Inches(2.3))
    earth.fill.solid()
    earth.fill.fore_color.rgb = rgb((0x0C, 0x24, 0x50))
    earth.line.color.rgb = rgb(ACCENT_BLUE)
    earth.line.width = Pt(1.2)
    # 经纬线（简单十字）
    add_line(s, gx - Inches(1.15), gy, gx + Inches(1.15), gy,
             color=ACCENT_BLUE, width=0.75)
    add_line(s, gx, gy - Inches(1.15), gx, gy + Inches(1.15),
             color=ACCENT_BLUE, width=0.75)
    # 节点圆点
    import math
    for ang_deg, c in [(30, GREEN), (110, ORANGE), (200, RED), (300, ACCENT_BLUE)]:
        rad = math.radians(ang_deg)
        dx = gx + Inches(0.9) * math.cos(rad)
        dy = gy + Inches(0.9) * math.sin(rad)
        nd = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                dx - Inches(0.07), dy - Inches(0.07),
                                Inches(0.14), Inches(0.14))
        nd.fill.solid()
        nd.fill.fore_color.rgb = rgb(c)
        nd.line.fill.background()

    # 对话示例
    dby = right.top + Inches(3.85)
    dialogs2 = [
        ('当前 CPU 的状态？', WHITE, False),
        ('当前 CPU 使用率 45.2%。', ACCENT_PURPLE, True),
        ('当前内存。', WHITE, False),
        ('当前内存 75.3%。', ACCENT_PURPLE, True),
    ]
    for txt, c, is_ai in dialogs2:
        off = Inches(0.2) if not is_ai else Inches(1.4)
        box = add_rect(s, right.left + off, dby,
                       right.width - off - Inches(0.4), Inches(0.38),
                       fill=BG_CARD, line=c, line_w=0.75, radius=True)
        add_text(s, box.left + Inches(0.12), box.top + Inches(0.06),
                 box.width - Inches(0.24), Inches(0.26),
                 ('👤 ' if not is_ai else '🤖 ') + txt,
                 size=10, color=c, bold=is_ai)
        dby += Inches(0.44)

    # 亮点
    hl = add_rect(s, right.left + Inches(0.25), dby + Inches(0.1),
                  right.width - Inches(0.5), Inches(0.6),
                  fill=BG_CARD, line=ORANGE, line_w=0.75, radius=True)
    add_text(s, hl.left + Inches(0.12), hl.top + Inches(0.08),
             hl.width - Inches(0.24), Inches(0.2),
             '✨ 关键亮点', size=10, bold=True, color=ORANGE)
    add_text(s, hl.left + Inches(0.12), hl.top + Inches(0.28),
             hl.width - Inches(0.24), Inches(0.3),
             'psutil 实时采集 · 百分比播报规范 · 3D 地球本地纹理',
             size=9, color=SILVER)


def slide_10_demo3_deploy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '功能演示 ③ · 排班 + 一键部署',
               '语音排班 · 二次确认 · 双击启动即运行',
               '08 · 排班 & 部署')
    slide_number(s, 10)

    # 左：排班
    left = add_rect(s, Inches(0.6), Inches(2.45), Inches(5.9), Inches(4.8),
                    fill=(0x08, 0x14, 0x2E), line=GREEN, line_w=1, radius=True)
    lb = add_rect(s, left.left + Inches(0.2), left.top + Inches(0.15),
                  Inches(1.8), Inches(0.35), fill=GREEN, radius=True)
    add_text(s, lb.left, lb.top + Inches(0.05), lb.width, Inches(0.25),
             '03 排班管理 + 语音排班', size=11, bold=True, color=BG_DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 3 场景
    scenes = [
        ('场景 1 · 语音排班',
         '👤 让王强排明天早班\n🤖 解析：王强 · 2026-08-05 · 早班',
         ACCENT_BLUE),
        ('场景 2 · 二次确认',
         '🤖 确认让王强排 2026-08-05 早班？\n👤 确认\n   ⟶ 更新 DB ⟶ "排班成功"',
         ORANGE),
        ('场景 3 · 日历视图',
         '7列×6行 月视图 · 四色班次\n早(蓝)/午(橙)/晚(紫)/休(灰)\n批量修改 · MERGE 策略',
         ACCENT_PURPLE),
    ]
    sy = left.top + Inches(0.62)
    for h, d, c in scenes:
        box = add_rect(s, left.left + Inches(0.2), sy,
                       left.width - Inches(0.4), Inches(1.2),
                       fill=BG_CARD, line=c, line_w=0.75, radius=True)
        add_text(s, box.left + Inches(0.15), box.top + Inches(0.08),
                 Inches(3), Inches(0.28), h, size=12, bold=True, color=c)
        add_text(s, box.left + Inches(0.15), box.top + Inches(0.4),
                 box.width - Inches(0.3), Inches(0.8), d,
                 size=10, color=SILVER)
        sy += Inches(1.3)

    # 右下小日历示意
    cal_y = left.top + Inches(4.4)
    cal_x = left.left + Inches(0.2)
    cw = left.width - Inches(0.4)
    ch = Inches(0.68)
    cal = add_rect(s, cal_x, cal_y, cw, ch,
                   fill=BG_CARD, line=GREEN, line_w=0.75, radius=True)
    # 表头
    for wk_i, wk in enumerate(['一', '二', '三', '四', '五', '六', '日']):
        add_text(s, cal_x + Inches(0.1) + wk_i * Inches(0.76),
                 cal.top + Inches(0.08), Inches(0.7), Inches(0.2),
                 wk, size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # 3 个班次色块
    slot_colors = [(3, ACCENT_BLUE, '早'), (5, ORANGE, '午'), (12, ACCENT_PURPLE, '晚')]
    for idx, c, nm in slot_colors:
        r, c_idx = idx // 7, idx % 7
        slot_x = cal_x + Inches(0.1) + c_idx * Inches(0.76)
        slot_y = cal.top + Inches(0.3) + r * Inches(0.16)
        tiny = add_rect(s, slot_x, slot_y, Inches(0.68), Inches(0.14),
                        fill=c, radius=True)
        add_text(s, tiny.left, tiny.top, tiny.width, tiny.height,
                 nm, size=7, bold=True, color=BG_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 右：一键部署 13 步
    right = add_rect(s, Inches(6.75), Inches(2.45), Inches(6.35), Inches(4.8),
                     fill=(0x08, 0x14, 0x2E), line=ACCENT_BLUE, line_w=1, radius=True)
    lb2 = add_rect(s, right.left + Inches(0.2), right.top + Inches(0.15),
                   Inches(3.4), Inches(0.35), fill=ACCENT_BLUE, radius=True)
    add_text(s, lb2.left, lb2.top + Inches(0.05), lb2.width, Inches(0.25),
             '🎛 一键部署 · 13 步全自动 (startup.bat)',
             size=11, bold=True, color=BG_DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # startup.bat 按钮示意
    bt = add_rect(s, right.left + Inches(1.5), right.top + Inches(0.62),
                  Inches(3.3), Inches(0.55),
                  fill=BG_CARD, line=GREEN, line_w=1.2, radius=True)
    add_text(s, bt.left, bt.top + Inches(0.12), bt.width, Inches(0.3),
             '▶  startup.bat   双击',
             size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 13 步两列编号
    steps = [
        ('1',  '检查 Python 3.12+'),
        ('2',  '检查 ffmpeg + PATH'),
        ('3',  '检查 Windows 神经网络 TTS'),
        ('4',  '创建 venv 虚拟环境'),
        ('5',  '安装后端 Python 依赖'),
        ('6',  '下载必需 TTS 模型 (zh-ll + melo)'),
        ('6b', '下载可选男声 (fanchen)'),
        ('7',  '检查 Node.js 22.18+'),
        ('8',  '安装前端 npm 依赖'),
        ('9',  '检查 Ollama 服务'),
        ('10', '拉取 qwen3:8b 模型'),
        ('11', '预加载模型常驻内存'),
        ('12', '从 OpsCenter.bak 还原数据库'),
        ('13', '启动后端(5000) + 前端(5173)'),
    ]
    sy = right.top + Inches(1.35)
    col_n = 2
    rows_per_col = 7
    for i, (num, desc) in enumerate(steps):
        col = i // rows_per_col
        row = i % rows_per_col
        cx = right.left + Inches(0.2) + col * Inches(3.0)
        cy = sy + row * Inches(0.52)
        c = ACCENT_BLUE if num != '6b' else ORANGE
        item = add_rect(s, cx, cy, Inches(2.95), Inches(0.42),
                        fill=BG_CARD, line=c, line_w=0.75, radius=True)
        # 编号圆
        nc = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                item.left + Inches(0.08), item.top + Inches(0.07),
                                Inches(0.28), Inches(0.28))
        nc.fill.solid()
        nc.fill.fore_color.rgb = rgb(c)
        nc.line.fill.background()
        add_text(s, nc.left, nc.top, nc.width, nc.height,
                 num, size=9, bold=True, color=BG_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, item.left + Inches(0.45), item.top + Inches(0.1),
                 Inches(2.4), Inches(0.22), desc,
                 size=9.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 底部下载策略小条
    ds = add_rect(s, right.left + Inches(0.2), right.top + Inches(6.35),
                  right.width - Inches(0.4), Inches(0.75),
                  fill=BG_CARD, line=ACCENT_PURPLE, line_w=0.75, radius=True)
    add_text(s, ds.left + Inches(0.15), ds.top + Inches(0.06),
             Inches(5.8), Inches(0.22),
             'TTS 模型下载智能策略',
             size=10, bold=True, color=ACCENT_PURPLE)
    add_text(s, ds.left + Inches(0.15), ds.top + Inches(0.3),
             Inches(5.8), Inches(0.4),
             '[SKIP] 已有则跳过  ·  [curl] 断点续传 10 次重试  ·  [WARN] 失败自动降级',
             size=9.5, color=SILVER)


def slide_11_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, '未来路线图',
               '从"听懂命令"到"主动预警" · 三阶段演进',
               '09 · 路线图')
    slide_number(s, 11)

    # 三阶段时间轴（横向）
    stages = [
        ('V1.0', '已完成 · 2026 Q3',
         ['语音对话 (ASR + TTS 全链路)',
          '11 类意图 LLM 解析 + 双重策略',
          '工单 + 统计看板 + 智能指派',
          '3D 监控 + psutil 实时采集',
          '排班管理 + 人员管理 + 岗位角色',
          '一键启动 + 全离线运行',
          'FSM 5 状态 + 二次确认护栏',
          'TTS 三引擎 + 多模型注册表'],
         GREEN, Inches(1.5)),
        ('V2.0', '规划中 · 2026 Q4 ~ 2027 Q1',
         ['屏幕理解 (截图 → 异常分析)',
          '知识库 RAG (运维文档智能检索)',
          'Docker 容器化部署',
          '移动端响应式适配',
          '多用户权限体系',
          '操作审计日志',
          '工单自动分类打标签'],
         ACCENT_PURPLE, Inches(5.8)),
        ('V3.0', '远景 · 2027 Q2+',
         ['自动故障诊断 (ML 异常检测)',
          '自动修复脚本 (审批后执行)',
          '预测性维护 (时序预测)',
          '告警自愈 (告警→诊断→修复闭环)',
          '多 Agent 协作 (工单/排班/监控)',
          '多模态输入 (语音+截图+日志)'],
         ORANGE, Inches(10.1)),
    ]
    # 中心时间轴线
    add_line(s, Inches(0.6), Inches(3.45), Inches(12.75), Inches(3.45),
             color=ACCENT_BLUE, width=3)
    for ver, subtitle, items, c, cx in stages:
        # 节点圆
        nd = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx - Inches(0.28), Inches(3.45) - Inches(0.28),
                                Inches(0.56), Inches(0.56))
        nd.fill.solid()
        nd.fill.fore_color.rgb = rgb(c)
        nd.line.color.rgb = rgb(BG_DEEP)
        nd.line.width = Pt(3)
        add_text(s, nd.left, nd.top + Inches(0.13), nd.width, Inches(0.3),
                 ver, size=12, bold=True, color=BG_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 版本标题卡片（上方）
        title_y = Inches(2.5)
        tcard = add_rect(s, cx - Inches(2.2), title_y, Inches(4.4), Inches(0.8),
                         fill=BG_CARD, line=c, line_w=1, radius=True)
        add_text(s, tcard.left, tcard.top + Inches(0.1),
                 tcard.width, Inches(0.3), f'{ver}  ·  {subtitle}',
                 size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        # 列表（下方）
        list_y = Inches(4.1)
        lcard = add_rect(s, cx - Inches(2.2), list_y, Inches(4.4), Inches(2.9),
                         fill=(0x08, 0x14, 0x2E), line=c, line_w=0.75, radius=True)
        # 勾选
        lines = []
        for it in items:
            lines.append(('[x] ' + it) if ver == 'V1.0' else ('[ ] ' + it))
        add_bullet(s, lcard.left + Inches(0.2), lcard.top + Inches(0.15),
                   lcard.width - Inches(0.4), lcard.height - Inches(0.3),
                   lines, size=10, color=SILVER)
        # 连接线
        add_line(s, cx, B(nd), cx, lcard.top, color=c, width=1.5)

    # 底部演进说明
    evolve = add_rect(s, Inches(0.6), Inches(6.98), Inches(12.5), Inches(0.4),
                      fill=(0x08, 0x14, 0x2E), line=ACCENT_BLUE, line_w=0.75, radius=True)
    add_text(s, evolve.left, evolve.top + Inches(0.08), evolve.width, Inches(0.24),
             '演进路径：智能运维助手  ──▶  多模态 Agent  ──▶  自主运维系统',
             size=12, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_12_end(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # 背景装饰圆环
    cx, cy = Inches(6.65), Inches(3.75)
    for r, col in [(Inches(3.2), (0x0A, 0x1F, 0x44)),
                   (Inches(2.3), (0x0F, 0x2B, 0x5E)),
                   (Inches(1.4), (0x1A, 0x3A, 0x7A))]:
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx - r, cy - r, r * 2, r * 2)
        ring.fill.background()
        ring.line.color.rgb = rgb(col)
        ring.line.width = Pt(2)
    core = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              cx - Inches(0.7), cy - Inches(0.7),
                              Inches(1.4), Inches(1.4))
    core.fill.solid()
    core.fill.fore_color.rgb = rgb(ACCENT_BLUE)
    core.line.fill.background()
    add_text(s, core.left, core.top + Inches(0.4),
             core.width, Inches(0.6), '🤖',
             size=34, color=BG_DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 标题
    add_text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.2),
             '轻言OPS', size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.55), Inches(12.1), Inches(0.5),
             '让运维工作"张口就来"',
             size=24, color=SILVER, align=PP_ALIGN.CENTER)

    # 三核心价值
    vals = [
        ('智能', '大模型 11 类意图 · 语音零门槛', ACCENT_BLUE),
        ('安全', '全离线 · FSM · 二次确认护栏', GREEN),
        ('便捷', '13 步一键 · 模型常驻', ORANGE),
    ]
    bx = Inches(0.6)
    by = Inches(5.5)
    bw = Inches(4.0)
    bh = Inches(0.95)
    for i, (t, d, c) in enumerate(vals):
        card = add_rect(s, bx + i * (bw + Inches(0.1)), by, bw, bh,
                        fill=BG_CARD, line=c, line_w=1, radius=True)
        add_text(s, card.left + Inches(0.2), card.top + Inches(0.12),
                 Inches(2.5), Inches(0.35), t, size=18, bold=True, color=c)
        add_text(s, card.left + Inches(0.2), card.top + Inches(0.55),
                 Inches(3.6), Inches(0.3), d, size=11, color=GRAY)

    # 五大技术亮点（底部横条）
    add_text(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.3),
             '五大技术亮点', size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    hl = ['全离线 AI 技术栈', '双重意图解析', 'FSM 5 状态护栏', '三引擎 TTS 降级', '智能数字处理']
    n = len(hl)
    g_w = Inches(12.1)
    gap = g_w / n
    for i, txt in enumerate(hl):
        tx = Inches(0.6) + gap * i
        card = add_rect(s, tx + Inches(0.05), Inches(6.95),
                        gap - Inches(0.1), Inches(0.42),
                        fill=(0x08, 0x14, 0x2E),
                        line=ACCENT_BLUE if i % 2 == 0 else ACCENT_PURPLE,
                        line_w=0.75, radius=True)
        c = ACCENT_BLUE if i % 2 == 0 else ACCENT_PURPLE
        add_text(s, card.left, card.top + Inches(0.1), card.width, Inches(0.24),
                 f'{i + 1}. {txt}',
                 size=10, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 联系方式底栏
    add_text(s, Inches(0.6), Inches(2.08), Inches(12.1), Inches(0.25),
             '闫吉乐  ·  2026 年 8 月',
             size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.35), Inches(12.1), Inches(0.25),
             'GitHub：github.com/YANJI-AFK/qingyan-OPS',
             size=11, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    # Q&A 大字（下方右）
    add_text(s, Inches(9.9), Inches(5.5), Inches(3.0), Inches(1.0),
             'Q & A', size=36, bold=True, color=ACCENT_PURPLE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# 主入口
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_contents(prs)
    slide_03_pain(prs)
    slide_04_position(prs)
    slide_05_ai_stack(prs)
    slide_06_fsm(prs)
    slide_07_voice(prs)
    slide_08_arch(prs)
    slide_09_demo12(prs)
    slide_10_demo3_deploy(prs)
    slide_11_roadmap(prs)
    slide_12_end(prs)

    prs.save(OUTPUT)
    print(f'✅ PPT 已生成: {OUTPUT}')
    print(f'   共 {len(prs.slides)} 页 · 16:9')


if __name__ == '__main__':
    try:
        main()
    except Exception as e:
        import traceback
        traceback.print_exc()
        sys.exit(1)
